from pptx import Presentation
from pptx.util import Cm, Pt, Inches
from pptx.enum.shapes import MSO_SHAPE
import pandas as pd
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_TICK_MARK, XL_TICK_LABEL_POSITION, XL_MARKER_STYLE
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import matplotlib.image as mpimg
from pptx.enum.action import PP_ACTION
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
import os

import time
# Record the start time
start_time = time.time()

def Hyperlink( run_object, source_slide, destination_slide ):
    rId = source_slide.part.relate_to(destination_slide.part, RT.SLIDE)
    rPr = run_object._r.get_or_add_rPr()
    hlinkClick = rPr.add_hlinkClick(rId)
    hlinkClick.set('action', 'ppaction://hlinksldjump')

project = "PT .."
gez = "Tropical Rainforest"
nicfi_date = " March 2023 "

# Specify the file name for the Excel file
table_path = '/.../ + project + '/'
excel_file = os.path.join(table_path, f"{project}_summary.xlsx")

map_folder = "/..."
map_list = pd.read_excel("/...")


# ### Page 1 - Title
prs = Presentation()

# Set slide width and height for A4 landscape
prs.slide_width = Inches(13.333)  # Width of A4 in inches (landscape)
prs.slide_height = Inches(7.5)  # Height of A4 in inches (landscape)

title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title


title.left, title.top, title.width, title.height = Cm(2), Cm(6), Cm(29.5), Cm(6) 

subtitle = slide.placeholders[1]
subtitle.left, subtitle.top, subtitle.width, subtitle.height = Cm(2), Cm(13), Cm(29.5), Cm(2) 


title.text = str( project +"CARBON PROJECT DESKTOP STUDY")
title.text_frame.paragraphs[0].font.size = Pt(54)  # Set font size in centimeters

subtitle.text = str(month + " 2024")
subtitle.text_frame.paragraphs[0].font.size = Pt(24)  # Set font size in centimeters


# ### Page 2 - Project Area
blank_slide_layout = prs.slide_layouts[5] # has no text box
slide = prs.slides.add_slide(blank_slide_layout)
shapes = slide.shapes

# Add title
slide_title = 'Project Area'

def add_title(slide_title):
    # Get title and body placeholders
    title_shape = shapes.title
    body_shape = shapes.placeholders[0]
    # Set text for title and body shapes
    title_shape.text = slide_title
    # Set font size, make title text bold, and align left
    title_shape.text_frame.paragraphs[0].font.size = Pt(30)  # Set font size in centimeters
    title_shape.text_frame.paragraphs[0].font.bold = True  # Make title text bold
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  # Align text to the left
    title_shape.text_frame.paragraphs[0].font.name = 'Calibri Light'
    title_shape.left, title_shape.top, title_shape.width, title_shape.height = Cm(1), Cm(1), Inches(10), Inches(0.76) 

add_title(slide_title)

# Add image
landmap_path = str(map_folder+map_list['Map'][0]+".jpg")

def image_layout(landmap_path):
    img_width = Inches(8.22)
    img_height = Inches(5.81)
    pic = slide.shapes.add_picture(landmap_path, left = Cm(0.8), top = Cm(3.5), width= img_width, height=img_height)
    
image_layout(landmap_path)

# Admin Table
Admin_data = pd.read_excel(excel_file, sheet_name="Admin")
Admin_data = round(Admin_data,3)

# Create a pivot table based on 'WADMKK', 'WADMKC', and 'WADMKD' with 'LUAS' as the value
project_area = Admin_data.pivot_table(index=['Kabupaten/Kota', 'Kecamatan', 'Desa'], values=['HA','%'], aggfunc='sum',margins=True)
project_area['HA'] = project_area['HA'].round(2)
project_area['%'] = project_area['%'].round(2)

project_area.reset_index(inplace=True)
project_area = project_area[['Kabupaten/Kota', 'Kecamatan', 'Desa', 'HA', '%']]
project_area['Kabupaten/Kota'][len(project_area)-1] = "TOTAL"


def table_layout(TableContent, font_size, bold):
    # Create a table on the slide
    left = Cm(22)
    top = Cm(3.5)  # Adjust the top position based on your layout
    width = Cm(11)
    height = Cm(1)
    table_placeholder = slide.shapes.add_table(rows=TableContent.shape[0] + 1, cols=TableContent.shape[1], left=left, top=top, width=width, height=height).table

    # Set table column names (header)
    for col_index, col_name in enumerate(TableContent.columns):
        cell = table_placeholder.cell(0, col_index)
        cell.text = col_name

        # Change font size for column names
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(font_size)  # Set the font size for column names

   # Populate table cells with data
    num_rows, num_cols = TableContent.shape  # Get the number of rows and columns
    for row in range(num_rows):
        for col in range(num_cols):
            table_placeholder.cell(row + 1, col).text = str(TableContent.iloc[row, col])


    # Populate table cells with data and set font style and type
    for row in range(TableContent.shape[0]):
        for col in range(TableContent.shape[1]):
            cell = table_placeholder.cell(row + 1, col)
            cell.text = str(TableContent.iloc[row, col])

            # Change font style and type for the cells
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
    #                 run.font.bold = True  # Set font to bold
    #                 run.font.italic = True  # Set font to italic
                    run.font.name = 'Calibri'  # Change font type (e.g., Arial)
                    run.font.size = Pt(font_size)
            
   # Access the last row of the table and set font to bold
    last_row_index = TableContent.shape[0]  # Get the index of the last row
    for col in range(TableContent.shape[1]):
        cell = table_placeholder.cell(last_row_index, col)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = bold  # Set font of the last row to bold
    
    
table_layout(TableContent = project_area, font_size = 8, bold = True)


# ### Page 3 - Area Ratio of Admin Within Concession
blank_slide_layout = prs.slide_layouts[5] # has no text box
slide = prs.slides.add_slide(blank_slide_layout)

shapes = slide.shapes

# Get title and body placeholders
title_shape = shapes.title
body_shape = shapes.placeholders[0]

# Add title
slide_title = 'Area Ratio of Admin Within Concession'
add_title(slide_title)

# Add image
adminmap_path = str(map_folder+map_list['Map'][1]+".jpg")
image_layout(adminmap_path)


# Admin Table
FullAdmin_data = pd.read_excel(excel_file, sheet_name="FullAdmin")
FullAdmin_data = round(FullAdmin_data,3)
total_row = pd.DataFrame({'Desa':['TOTAL'],
                  'Luas Desa Dalam Konsesi (HA)':sum(FullAdmin_data['Luas Desa Dalam Konsesi (HA)']),
                  'Luas Desa Dalam Konsesi (%)': sum(FullAdmin_data['Luas Desa Dalam Konsesi (%)']),
                  'Total Luas Desa (HA)':'-',
                  'Persentase Total Luas Desa yg Masuk Kedalam Konsesi (%)':'-'})

FullAdmin_data2 = pd.concat([FullAdmin_data, total_row], ignore_index=True)
FullAdmin_data2

table_layout(TableContent = FullAdmin_data2, font_size = 8, bold = True)



# ### Page 4 - Land Cover
blank_slide_layout = prs.slide_layouts[5] # has no text box
slide = prs.slides.add_slide(blank_slide_layout)

shapes = slide.shapes

# Get title and body placeholders
title_shape = shapes.title
body_shape = shapes.placeholders[0]


# Add title
slide_title = 'Land Cover'
add_title(slide_title)

# Add image
landmap_path = str(map_folder+map_list['Map'][2]+".jpg")
image_layout(landmap_path)


# Admin Table
PL_data = pd.read_excel(excel_file, sheet_name="PL")
PL_data

table_layout(TableContent = PL_data, font_size = 11, bold = True)


# ### Page 5 - Forest Map
blank_slide_layout = prs.slide_layouts[5] # has no text box
slide = prs.slides.add_slide(blank_slide_layout)

shapes = slide.shapes

# Get title and body placeholders
title_shape = shapes.title
body_shape = shapes.placeholders[0]

# Add title
slide_title = 'Forest Map'
add_title(slide_title)

# Add image
forestmap_path = str(map_folder+map_list['Map'][3]+".jpg")
image_layout(forestmap_path)

# Admin Table
Hutan_data = pd.read_excel(excel_file, sheet_name="Hutan")
Hutan_data


table_layout(TableContent = Hutan_data, font_size = 11, bold = True)


# ### Page 6 - ESA World Cover
blank_slide_layout = prs.slide_layouts[5] # has no text box
slide = prs.slides.add_slide(blank_slide_layout)

shapes = slide.shapes

# Get title and body placeholders
title_shape = shapes.title
body_shape = shapes.placeholders[0]

# Add title
slide_title = 'ESA World Cover'
add_title(slide_title)

# Add image
WorldCovermap_path = str(map_folder+map_list['Map'][4]+".jpg")
image_layout(WorldCovermap_path)

# Admin Table
WorldCover_data = pd.read_excel(excel_file, sheet_name="WorldCover")
table_layout(TableContent = WorldCover_data, font_size = 11, bold = True)



# ### Page 7 - Fire Case
blank_slide_layout = prs.slide_layouts[5] # has no text box
slide = prs.slides.add_slide(blank_slide_layout)

shapes = slide.shapes

# Get title and body placeholders
title_shape = shapes.title
body_shape = shapes.placeholders[0]

# Add title
slide_title = 'Fire Case'
add_title(slide_title)

# Add image
firemap_path = str(map_folder+map_list['Map'][5]+".jpg")
image_layout(firemap_path)

# Text Book
txBox = slide.shapes.add_textbox(left = Cm(22), top = Cm(3.1), width = Cm(11), height = Cm(11))

# FCL Table
Fire = pd.read_excel(excel_file, sheet_name="Fire")
Fire['HA'] = round(Fire['HA'],2)

# Extract the 'Year' and 'HA' columns for the chart
years = Fire['Year'].tolist()
ha_values = Fire['HA'].tolist()

# Define chart data
chart_data = CategoryChartData()
chart_data.categories = years  # X-axis data (Years)
chart_data.add_series('Hectares', ha_values)  # Y-axis data (HA)

# Define chart size and position
x, y, cx, cy = Cm(22), Cm(3.1), Cm(11), Inches(4)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
).chart

# Set chart title and font style/size
chart.has_title = True
chart.chart_title.text_frame.text = 'Annual Forest Loss Due to Fire (HA)'
title_format = chart.chart_title.text_frame.paragraphs[0].font
title_format.name = 'Calibri'
title_format.size = Pt(16)  # Font size in inches (adjust as needed)
title_format.bold = True  # Adjust other properties as needed

# Set font style/size for X-axis labels
category_axis = chart.category_axis
category_axis.tick_labels.font.size = Pt(8)  # Font size in inches for X-axis labels

# Set font style/size for Y-axis labels
value_axis = chart.value_axis
value_axis.tick_labels.font.size = Pt(8)  # Font size in inches for Y-axis labels

# Change bar color to red and add data labels
plot = chart.plots[0]
for index, point in enumerate(plot.series[0].points):
    fill = point.format.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(237, 125, 49)  # Red color for bars (RGB values as a tuple)
    
   # Add data labels
    data_label = point.data_label
    data_label.text_frame.text = str(ha_values[index])  # Display value as label
    for paragraph in data_label.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(8)  # Adjust font size for data labels
            
# Gridlines
chart.category_axis.has_major_gridlines = False #Remove Vertical Major Gridlines
chart.value_axis.has_major_gridlines = False #Remove Horizontal major Gridlines



# ### Page 8 - Forest Cover Lost
blank_slide_layout = prs.slide_layouts[5] # has no text box
slide = prs.slides.add_slide(blank_slide_layout)

shapes = slide.shapes

# Get title and body placeholders
title_shape = shapes.title
body_shape = shapes.placeholders[0]


# Add title
slide_title = 'Forest Cover Lost'
add_title(slide_title)

# Add image
firemap_path = str(map_folder+map_list['Map'][6]+".jpg")
image_layout(firemap_path)


# Text Book
txBox = slide.shapes.add_textbox(left = Cm(22), top = Cm(2.5), width = Cm(11), height = Cm(7))
tf = txBox.text_frame

# FCL Table
FCL = pd.read_excel(excel_file, sheet_name="FCL")
FCL['HA'] = round(FCL['HA'], 2)

total_kehilangan_hutan = round(sum(FCL['HA']),2)
presentase_kehilangan_hutan = round(total_kehilangan_hutan/project_area['HA'].iloc[0:len(project_area)-1].sum()*100,0) 
kehilangan_max = max(FCL['HA'])
tahun_kehilangan_max = FCL.loc[FCL['HA'].idxmax(), 'Year']
Fire_total = round(sum(Fire['HA']),2)

# Set the text with line breaks
text_with_line_breaks = (
    "Sejak tahun 2001 hingga tahun 2022 telah \n"
    "terjadi total " + str(total_kehilangan_hutan) +" ha kehilangan tutupan\n"
    "hutan, sekitar " + str(presentase_kehilangan_hutan) +"% dari total kawasan.\n"
    "\n"
    "Kehilangan hutan Tertinggi terjadi pada \n"
    "tahun " + str(tahun_kehilangan_max) + " dengan kehilangan hutan \n"
    "sebesar " + str(kehilangan_max) + " ha.\n"
    "\n"
    ""+ str(Fire_total) +" ha hutan hilang diakibatkan oleh\n" 
    "kebakaran hutan." 
)

# Set the text within the text box
lines = text_with_line_breaks.split('\n')
for line in lines:
    p = tf.add_paragraph()
    p.text = line
    # Adjust line spacing for compactness (modify this value as needed)
    p.line_spacing = Pt(18)  # Adjust this value to change line spacing


# Extract the 'Year' and 'HA' columns for the chart
years = FCL['Year'].tolist()
ha_values = FCL['HA'].tolist()

# Define chart data
chart_data = CategoryChartData()
chart_data.categories = years  # X-axis data (Years)
chart_data.add_series('Hectares', ha_values)  # Y-axis data (HA)

# Define chart size and position
x, y, cx, cy = Cm(22), Cm(10), Cm(11), Inches(3.3)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data  # Change chart type to LINE
).chart

# Remove legend
chart.has_legend = False

# Set chart title and font style/size
chart.has_title = True
chart.chart_title.text_frame.text = 'Forest Cover Lost (HA)'
title_format = chart.chart_title.text_frame.paragraphs[0].font
title_format.name = 'Calibri'
title_format.size = Pt(16)  # Font size in points (adjust as needed)
title_format.bold = True  # Adjust other properties as needed

# Set font style/size for X-axis labels
category_axis = chart.category_axis
category_axis.tick_labels.font.size = Pt(8)  # Font size in points for X-axis labels

# Set font style/size for Y-axis labels
value_axis = chart.value_axis
value_axis.tick_labels.font.size = Pt(8)  # Font size in points for Y-axis labels

            
# Add data labels to the line chart
plot = chart.plots[0]
for series in plot.series:
    series.has_data_labels = True
    for point, value in zip(series.points, ha_values):
        data_label = point.data_label
        data_label.text = str(value)
        data_label.font.size = Pt(8)  # Font size for data labels

# Access the first series (change the index if needed)
line_series = chart.series[0]

# Change the line color to red (RGB: 255, 0, 0)
line = line_series.format.line
line.color.rgb = RGBColor(255, 0, 0)  # Red color

# Optional: Change line width
line.width = Pt(2)  # Set line width in points (adjust as needed)

# Gridlines
chart.category_axis.has_major_gridlines = False #Remove Vertical Major Gridlines
chart.value_axis.has_major_gridlines = False #Remove Horizontal major Gridlines

# Iterate through each series and set marker properties
for series in chart.series:
    series.marker.style = XL_MARKER_STYLE.CIRCLE  # Marker style: CIRCLE, SQUARE, DIAMOND, etc.
    series.marker.size = 5  # Marker size (adjust as needed)
    # Customize marker fill color (optional)
    series.marker.format.fill.solid()
    # series.marker.format.fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red color



# ### Page 9 - Carbon Stock
blank_slide_layout = prs.slide_layouts[5] # has no text box
slide = prs.slides.add_slide(blank_slide_layout)

shapes = slide.shapes

# Get title and body placeholders
title_shape = shapes.title
body_shape = shapes.placeholders[0]

# Add title
slide_title = 'Carbon Stock'
add_title(slide_title)

# Add image
firemap_path = str(map_folder+map_list['Map'][7]+".jpg")
image_layout(firemap_path)

# Text Book
txBox = slide.shapes.add_textbox(left = Cm(22), top = Cm(2.5), width = Cm(11), height = Cm(11))
tf = txBox.text_frame

# Set the text with line breaks
text_with_line_breaks = (
    "AGB (Aboveground Biomass)\n"
    "menggunakan data dari GlobBiomass.\n"
    "\n"
    "Perhitungan BGB (Belowground Biomass)\n"
    "menggunakan default root shoot ratio.\n"
    "Dimana kawasan ini masuk pada Global \n" 
    "Ecological Zones " + str(gez) +  " dari \n" 
    "IPCC 2019 Refinement to the 2006.\n"
    "\n"
    "CO2 dalah penangkapan dan\n"
    "penyimpanan karbon dioksida dari \n"
    "atmosfer dalam jangka waktu yang lama.\n" 
    "Dihitung dengan Total Carbon Stock*3.67 \n" 
    "(Koefisien Berat Carbon)." 
)

# Set the text within the text box
lines = text_with_line_breaks.split('\n')
for line in lines:
    p = tf.add_paragraph()
    p.text = line
    # Adjust line spacing for compactness (modify this value as needed)
    p.line_spacing = Pt(22)  # Adjust this value to change line spacing

# Biomass Table
biomass_data = pd.read_excel(excel_file, sheet_name="Biomass")
biomass_data = round(biomass_data, 2)

# Create a table on the slide
left = Cm(22)
top = Cm(15.3)
width = Cm(11)
height = Cm(1.8)
table_placeholder = slide.shapes.add_table(rows=biomass_data.shape[0] + 1, cols=biomass_data.shape[1], left=left, top=top, width=width, height=height).table


# Set table column names (header)
for col_index, col_name in enumerate(biomass_data.columns):
    cell = table_placeholder.cell(0, col_index)
    cell.text = col_name

    # Change font size for column names
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(11)  # Set the font size for column names


# Populate table cells with data and set font size
for row in range(biomass_data.shape[0]):
    for col in range(biomass_data.shape[1]):
        cell = table_placeholder.cell(row + 1, col)
        cell.text = str(biomass_data.iloc[row, col])

        # Change font size for the cells
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(11)  # Set font size (e.g., 14 points)




# ### Page 10 - Forest Height
blank_slide_layout = prs.slide_layouts[5] # has no text box
slide = prs.slides.add_slide(blank_slide_layout)

shapes = slide.shapes

# Get title and body placeholders
title_shape = shapes.title
body_shape = shapes.placeholders[0]

# Add title
slide_title = 'Forest Height'
add_title(slide_title)

# Add image
chmmap_path = str(map_folder+map_list['Map'][8]+".jpg")
image_layout(chmmap_path)

# Text Book
txBox = slide.shapes.add_textbox(left = Cm(22), top = Cm(2.5), width = Cm(11), height = Cm(11))

tf = txBox.text_frame

# Set the text with line breaks
text_with_line_breaks = (
    "Mewakili Ketinggian Hutan tahun 2020 \n"
    "dalam resolusi 10m.\n"
    "\n"
    "Dikembangkan menggunakan model\n"
    "Deep Learning dari kombinasi GEDI \n"
    "(satelit LiDAR) & Sentinel-2 (satelit Multi- \n" 
    "spektral)" 
)

# Set the text within the text box
lines = text_with_line_breaks.split('\n')
for line in lines:
    p = tf.add_paragraph()
    p.text = line
    # Adjust line spacing for compactness (modify this value as needed)
    p.line_spacing = Pt(22)  # Adjust this value to change line spacing

# chm Table
chm_data = pd.read_excel(excel_file, sheet_name="CHM")
chm_data = round(chm_data, 2)

# Create a table on the slide
left = Cm(22)
top = Cm(9.5)
width = Cm(11)
height = Cm(1.8)
table_placeholder = slide.shapes.add_table(rows=chm_data.shape[0] + 1, cols=chm_data.shape[1], left=left, top=top, width=width, height=height).table


# Set table column names (header)
for col_index, col_name in enumerate(chm_data.columns):
    cell = table_placeholder.cell(0, col_index)
    cell.text = col_name

    # Change font size for column names
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(11)  # Set the font size for column names


# Populate table cells with data and set font size
for row in range(chm_data.shape[0]):
    for col in range(chm_data.shape[1]):
        cell = table_placeholder.cell(row + 1, col)
        cell.text = str(chm_data.iloc[row, col])

        # Change font size for the cells
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(11)  # Set font size (e.g., 14 points)
              
   # Access the last row of the table and set font to bold
    last_row_index = chm_data.shape[0]  # Get the index of the last row
    for col in range(chm_data.shape[1]):
        cell = table_placeholder.cell(last_row_index, col)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True  # Set font of the last row to bold




# ### Page 11 - Mining Area
blank_slide_layout = prs.slide_layouts[5] # has no text box
slide = prs.slides.add_slide(blank_slide_layout)

shapes = slide.shapes

# Get title and body placeholders
title_shape = shapes.title
body_shape = shapes.placeholders[0]

# Add title
slide_title = 'Mining Area'
add_title(slide_title)

# Add image
miningemap_path = str(map_folder+map_list['Map'][9]+".jpg")
image_layout(miningemap_path)

# Add Table
mining_data = pd.read_excel(excel_file, sheet_name="Mining")
table_layout(TableContent = mining_data, font_size = 10, bold = False)




# ### Page 12 - Concession Map
blank_slide_layout = prs.slide_layouts[5] # has no text box
slide = prs.slides.add_slide(blank_slide_layout)

shapes = slide.shapes

# Get title and body placeholders
title_shape = shapes.title
body_shape = shapes.placeholders[0]

# Add title
slide_title = 'Concession Map'
add_title(slide_title)

# Add image
concessionmap_path = str(map_folder+map_list['Map'][10]+".jpg")
image_layout(concessionmap_path)



# ### Page 13 - Land System
blank_slide_layout = prs.slide_layouts[5] # has no text box
slide = prs.slides.add_slide(blank_slide_layout)

shapes = slide.shapes

# Get title and body placeholders
title_shape = shapes.title
body_shape = shapes.placeholders[0]


# Add title
slide_title = 'Land System'
add_title(slide_title)

# Add image
landsystemmap_path = str(map_folder+map_list['Map'][11]+".jpg")
image_layout(landsystemmap_path)

# Add Table
LS_data = pd.read_excel(excel_file, sheet_name="LandSystem")
# remove Nan Table
LS_data.iloc[-1, 1] = "" # -1 represents the last row, 1 represents the second column (index starts from 0)
LS_data

def table_layout_merged(TableContent, font_size, bold):
    # Create a table on the slide
    left = Cm(22)
    top = Cm(3.5)  # Adjust the top position based on your layout
    width = Cm(11)
    height = Cm(1)
    table_placeholder = slide.shapes.add_table(rows=TableContent.shape[0] + 1, cols=TableContent.shape[1], left=left, top=top, width=width, height=height).table

    # Set table column names (header)
    for col_index, col_name in enumerate(TableContent.columns):
        cell = table_placeholder.cell(0, col_index)
        cell.text = col_name

        # Change font size for column names
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(font_size)  # Set the font size for column names

   # Populate table cells with data
    num_rows, num_cols = TableContent.shape  # Get the number of rows and columns
    for row in range(num_rows):
        for col in range(num_cols):
            table_placeholder.cell(row + 1, col).text = str(TableContent.iloc[row, col])


    # Populate table cells with data and set font style and type
    for row in range(TableContent.shape[0]):
        for col in range(TableContent.shape[1]):
            cell = table_placeholder.cell(row + 1, col)
            cell.text = str(TableContent.iloc[row, col])

            # Change font style and type for the cells
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    # run.font.bold = True  # Set font to bold
                    # run.font.italic = True  # Set font to italic
                    run.font.name = 'Calibri'  # Change font type (e.g., Arial)
                    run.font.size = Pt(font_size)
            
   # Access the last row of the table and set font to bold
    last_row_index = TableContent.shape[0]  # Get the index of the last row
    for col in range(TableContent.shape[1]):
        cell = table_placeholder.cell(last_row_index, col)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = bold  # Set font of the last row to bold
                
    # Access the last row of the table and set font to bold
    last_row_index = TableContent.shape[0]  # Get the index of the last row
    for col in range(TableContent.shape[1]):
        cell = table_placeholder.cell(last_row_index, col)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = bold  # Set font of the last row to bold

    # Merge the first and second columns of the last row
    merged_cell = table_placeholder.cell(last_row_index, 0).merge(table_placeholder.cell(last_row_index, 1))
    #merged_cell.text = 'TOTAL'  # Set the text for the merged cell
    
    return table_placeholder  # Return the modified table object

                

table_layout_merged(TableContent = LS_data, font_size = 11, bold = True)
    


# ### Page 14 - Soil Map
blank_slide_layout = prs.slide_layouts[5] # has no text box
slide = prs.slides.add_slide(blank_slide_layout)

shapes = slide.shapes

# Get title and body placeholders
title_shape = shapes.title
body_shape = shapes.placeholders[0]


# Add title
slide_title = 'Soil Map'
add_title(slide_title)

# Add image
soilmap_path = str(map_folder+map_list['Map'][12]+".jpg")
image_layout(soilmap_path)

# Add table
soil_data = pd.read_excel(excel_file, sheet_name="Soil")
# remove Nan Table
soil_data.iloc[-1, 1] = "" # -1 represents the last row, 1 represents the second column (index starts from 0)
soil_data
table_layout_merged(TableContent = soil_data, font_size = 11, bold = True)



# ### Page 15 - Rainfall
rainfall_slide_layout = prs.slide_layouts[5] # has no text box
rainfall_slide = prs.slides.add_slide(rainfall_slide_layout)

shapes = rainfall_slide.shapes

# Get title and body placeholders
title_shape = shapes.title
body_shape = shapes.placeholders[0]

# Add title
slide_title = 'Rainfall'
add_title(slide_title)

# # Add image
Rainfallmap_path = str(map_folder+map_list['Map'][13]+".jpg")
# image_layout(Rainfallmap_path)

# Add image
img_width = Inches(8.22)
img_height = Inches(5.81)
pic = rainfall_slide.shapes.add_picture(Rainfallmap_path, left = Cm(0.8), top = Cm(3.5), width= img_width, height=img_height)


# Text Book
txBox = rainfall_slide.shapes.add_textbox(left = Cm(22), top = Cm(2.5), width = Cm(11), height = Cm(11))
tf = txBox.text_frame


# Extract rainfall data
rainfall_data = pd.read_excel(excel_file, sheet_name="Rainfall")

# Filter the DataFrame to extract 'Total Rainfall' where 'Year' contains 'Yearly Average'
annual_rainfal = rainfall_data.query("Year=='Yearly Average'")["Total Rainfall"].iloc[0]
annual_rainfal


# Set the text with line breaks
text_with_line_breaks = (
    "CHIRPS adalah database curah hujan \n"
    "daratan yang merupakan kombinasi dari \n"
    "tiga informasi curah hujan yaitu \n"
    "klimatologi global, estimasi curah hujan \n"
    "berbasis satelit, dan curah hujan hasil \n"
    "pengamatan in-situ\n" 
    "\n"
    "Area kajian memiliki jumlah hujan \n"
    "tahunan sebesar " + str(annual_rainfal) + " mm/th \n"
    "(berdasarkan data tahun 2010-2022)"
)

# Set the text within the text box
lines = text_with_line_breaks.split('\n')
for line in lines:
    p = tf.add_paragraph()
    p.text = line
    # Adjust line spacing for compactness (modify this value as needed)
    p.line_spacing = Pt(22)  # Adjust this value to change line spacing


shape = shapes.add_shape(MSO_SHAPE.PENTAGON, left = Cm(28) , top = Cm(17.4), width = Inches(1.75), height = Inches(0.31))
shape.text = 'Rainfall Chart'

# Change the shape fill color
fill = shape.fill
fill.solid()
fill.fore_color.rgb = RGBColor(225, 225, 0)  # Change the fill color to a shade of green (RGB values as a tuple)

# Access the text frame and set the font size
text_frame = shape.text_frame ## change into text_frame as a hyperlink
for paragraph in text_frame.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(12)  # Change the font size to 18 points (adjust as needed)
        # run.font.color.rgb = RGBColor(237, 0, 0)
        
left = left + width - Inches(0.4)
width = Inches(2.0)  # chevrons need more width for visual balance



# ### Page 16 - Protected Areas
blank_slide_layout = prs.slide_layouts[5] # has no text box
slide = prs.slides.add_slide(blank_slide_layout)

shapes = slide.shapes

# Get title and body placeholders
title_shape = shapes.title
body_shape = shapes.placeholders[0]

# Add title
slide_title = 'Protected Area'
add_title(slide_title)

# Add image
protectedareasmap_path = str(map_folder+map_list['Map'][14]+".jpg")
image_layout(protectedareasmap_path)


# ### Page 17 - Peatland Map
blank_slide_layout = prs.slide_layouts[5] # has no text box
slide = prs.slides.add_slide(blank_slide_layout)

shapes = slide.shapes

# Get title and body placeholders
title_shape = shapes.title
body_shape = shapes.placeholders[0]

# Add title
slide_title = 'Peatland'
add_title(slide_title)

# Add image
peatlandmap_path = str(map_folder+map_list['Map'][15]+".jpg")
image_layout(peatlandmap_path)
# Add Table
peat_data = pd.read_excel(excel_file, sheet_name="Peatland")
table_layout(TableContent = peat_data, font_size = 11, bold = True)


# ### Page 18 - Wetland
blank_slide_layout = prs.slide_layouts[5] # has no text box
slide = prs.slides.add_slide(blank_slide_layout)

shapes = slide.shapes

# Get title and body placeholders
title_shape = shapes.title
body_shape = shapes.placeholders[0]

# Add title
slide_title = 'Wetlands'
add_title(slide_title)

# Add image
peatlandmap_path = str(map_folder+map_list['Map'][16]+".jpg")
image_layout(peatlandmap_path)

# Add Table
wetland_data = pd.read_excel(excel_file, sheet_name="Wetlands")
table_layout(TableContent = wetland_data, font_size = 11, bold = True)



# ### Page 19 - Mangrove
blank_slide_layout = prs.slide_layouts[5] # has no text box
slide = prs.slides.add_slide(blank_slide_layout)

shapes = slide.shapes

# Get title and body placeholders
title_shape = shapes.title
body_shape = shapes.placeholders[0]

# Add title
slide_title = 'Mangrove'
add_title(slide_title)

# Add image
peatlandmap_path = str(map_folder+map_list['Map'][17]+".jpg")
image_layout(peatlandmap_path)
# Add Table
mangrove_data = pd.read_excel(excel_file, sheet_name="Mangrove")
table_layout(TableContent = mangrove_data, font_size = 11, bold = True)



# ### Page 20 - Elevation Map
blank_slide_layout = prs.slide_layouts[5] # has no text box
slide = prs.slides.add_slide(blank_slide_layout)

shapes = slide.shapes

# Get title and body placeholders
title_shape = shapes.title
body_shape = shapes.placeholders[0]

# Add title
slide_title = 'Elevation Map'
add_title(slide_title)

# Add image
elevationmap_path = str(map_folder+map_list['Map'][18]+".jpg")
def img_middle(img):
    img_width = Inches(8.22)
    img_height = Inches(5.81)
    pic = slide.shapes.add_picture(img, left = Cm(6.5), top = Cm(3.5), width= img_width, height=img_height)

img_middle(elevationmap_path)


# ### Page 21 - Slope Map
blank_slide_layout = prs.slide_layouts[5] # has no text box
slide = prs.slides.add_slide(blank_slide_layout)

shapes = slide.shapes

# Get title and body placeholders
title_shape = shapes.title
body_shape = shapes.placeholders[0]

# Add title
slide_title = 'Slope Map'
add_title(slide_title)

# Add image
slopemap_path = str(map_folder+map_list['Map'][19]+".jpg")
def img_middle(img):
    img_width = Inches(8.22)
    img_height = Inches(5.81)
    pic = slide.shapes.add_picture(img, left = Cm(6.5), top = Cm(3.5), width= img_width, height=img_height)

img_middle(slopemap_path)



# ### Page 22 - Appendix
blank_slide_layout = prs.slide_layouts[5] # has no text box
slide = prs.slides.add_slide(blank_slide_layout)

shapes = slide.shapes

# Get title and body placeholders
title_shape = shapes.title
body_shape = shapes.placeholders[0]

# Add title
slide_title = 'Appendix'
add_title(slide_title)

# Add image
satellitemap_path = str(map_folder+map_list['Map'][20]+".jpg")
def img_middle(img):
    img_width = Inches(8.22)
    img_height = Inches(5.81)
    pic = slide.shapes.add_picture(img, left = Cm(6.5), top = Cm(3.5), width= img_width, height=img_height)

img_middle(satellitemap_path)



# ### Page 23 NICFI Imagery
blank_slide_layout = prs.slide_layouts[5] # has no text box
slide = prs.slides.add_slide(blank_slide_layout)

shapes = slide.shapes

# Get title and body placeholders
title_shape = shapes.title
body_shape = shapes.placeholders[0]

# Add title
slide_title = "NICFI Imagery " + str(nicfi_date) + " (4.77m Resolution)"
add_title(slide_title)

# Add image
satellitemap_path = str(map_folder+map_list['Map'][21]+".jpg")
def img_middle(img):
    img_width = Inches(8.22)
    img_height = Inches(5.81)
    pic = slide.shapes.add_picture(img, left = Cm(6.5), top = Cm(3.5), width= img_width, height=img_height)

img_middle(satellitemap_path)




# ### Page 24 - Rainfall Chart
blank_slide_layout = prs.slide_layouts[5] # has no text box
slide = prs.slides.add_slide(blank_slide_layout)

shapes = slide.shapes

# Get title and body placeholders
title_shape = shapes.title
body_shape = shapes.placeholders[0]

# Add title
slide_title = "Rainfall Chart"
add_title(slide_title)


def rainfall_chart(sheet_name, top, rainfall_title):
    # FCL Table
    rainfall = pd.read_excel(excel_file, sheet_name= sheet_name)

    # Extract the 'Year' and 'HA' columns for the chart
    date = rainfall['Month'].tolist()
    values = rainfall['Curah Hujan'].tolist()

    # Define chart data
    chart_data = CategoryChartData()
    chart_data.categories = date  # X-axis data (Years)
    chart_data.add_series('Curah Hujan', values)  # Y-axis data (HA)

    # Define chart size and position
    x, y, cx, cy = Cm(1), Cm(top), Cm(32), Cm(7)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data  # Change chart type to LINE
    ).chart

    # Remove legend
    chart.has_legend = False

    # Set chart title and font style/size
    chart.has_title = True
    chart.chart_title.text_frame.text = rainfall_title
    title_format = chart.chart_title.text_frame.paragraphs[0].font
    title_format.name = 'Calibri'
    title_format.size = Pt(14)  # Font size in points (adjust as needed)
    title_format.bold = True  # Adjust other properties as needed

    # Set font style/size for X-axis labels
    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(9)  # Font size in points for X-axis labels
    category_axis.tick_labels.rotation = 45  # Set the rotation angle (e.g., 45 degrees)

    # Set font style/size for Y-axis labels
    value_axis = chart.value_axis
    value_axis.tick_labels.font.size = Pt(9)  # Font size in points for Y-axis labels


rainfall_chart("Rainfall-Sum", 3.5, "Jumlah Curah Hujan Bulanan 2010-2022 (mm/bln)")
rainfall_chart("Rainfall-Mean", 11, "Rerata Curah Bulanan 2010-2022 (mm/bln)")


# Create the hyperlink from the text box on slide 2 to the last slide ('Rainfall Chart')
run_object = text_frame.paragraphs[0].runs[0]  # Assume this is the run object where you want to apply the hyperlink
Hyperlink(run_object, rainfall_slide, slide)

# Calculate elapsed time
end_time = time.time()
elapsed_time = (end_time - start_time)#/60

print("Elapsed time:", elapsed_time, "second")


prs.save('/...pptx')


